
from pptx import Presentation
from pptx.util import Inches, Pt
import re


# Input/Output paths
SLIDES_MD_PATH = r"h:\マイドライブ\antigravity\scratch\Business_Plans\03_APE\00_事業計画資料_Docs\Slides.md"
OUTPUT_PPTX_PATH = r"h:\マイドライブ\antigravity\scratch\Business_Plans\03_APE\00_事業計画資料_Docs\APE_Presentation.pptx"

def parse_slides(md_path):
    """Parses a markdown file into a list of slide dicts."""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by "## Slide" or "---" as dividers
    # Assuming each major section starts with "## Slide X: Title" or just "## Title"
    raw_slides = re.split(r'\n---\n', content)
    
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        
        lines = raw.split('\n')
        title = ""
        body_lines = []
        
        for line in lines:
            # Find title (## Slide X: Title or ## Title)
            if line.startswith('## '):
                # Remove "## Slide X: " prefix if present
                title = re.sub(r'^## (Slide \d+:\s*)?', '', line).strip()
            elif line.startswith('# '):
                title = line[2:].strip()
            else:
                body_lines.append(line)
        
        slides.append({'title': title, 'body': '\n'.join(body_lines).strip()})
    
    return slides

def create_pptx(slides_data, output_path):
    """Creates a PPTX from parsed slide data."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    title_slide_layout = prs.slide_layouts[5] # Blank
    
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Add Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data['title']
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        
        # Add Body
        body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(7))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        
        body_lines = slide_data['body'].split('\n')
        for line in body_lines:
            p = body_frame.add_paragraph()
            # Basic bullet handling
            if line.strip().startswith('* ') or line.strip().startswith('- '):
                p.text = '  • ' + line.strip()[2:]
                p.level = 0
            elif line.strip().startswith('    * ') or line.strip().startswith('    - '):
                p.text = '      ◦ ' + line.strip()[6:]
                p.level = 1
            else:
                p.text = line
            p.font.size = Pt(18)
        
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    slides = parse_slides(SLIDES_MD_PATH)
    print(f"Parsed {len(slides)} slides.")
    create_pptx(slides, OUTPUT_PPTX_PATH)
